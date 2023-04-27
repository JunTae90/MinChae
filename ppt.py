from PySide6.QtCore import *
import re
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.dml.color import *
from PIL import Image
class Maker(QThread):
    start = Signal()
    end = Signal()
    def __init__(self, parent=None):
        super().__init__(parent)
        self.imgList = []
    def run(self):
        self.start.emit()
        width = 25.4
        height = 19.05
        max_width = 25.4/3
        max_height = 19.05/2
        prs = Presentation()
        prs.slide_width = Cm(width)
        prs.slide_height = Cm(height)
        slide_layout = prs.slide_layouts[6]

        for key in self.imgNames.keys():
            imgNames = self.imgNames[key]
            if len(imgNames) == 0:
                continue
            slide = prs.slides.add_slide(slide_layout)
            tb = slide.shapes.add_textbox(Cm(0), Cm(0), Cm(width/10), Cm(height/20))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = key
            font = run.font
            font.color.rgb = RGBColor(255, 0, 0)
            for i in range(0, len(imgNames)):
                if i // 6 > 0 and i % 6 == 0:
                    slide = prs.slides.add_slide(slide_layout)
                im = Image.open(self.path+imgNames[i])
                x = max_width*((i % 6) % 3)
                y = max_height*((i % 6) // 3)
                if im.size[0] > im.size[1]:
                    pc = slide.shapes.add_picture(self.path+imgNames[i], Cm(x), Cm(y), width = Cm(max_width))
                else:
                    pc = slide.shapes.add_picture(self.path + imgNames[i], Cm(x), Cm(y), height=Cm(max_height))
                slide.shapes._spTree.insert(2, pc._element)

        prs.save(self.path+"demo.pptx")
        self.end.emit()
    @Slot(str, result=dict)
    def extractInfo(self, chatPath):
        res = {}
        cur = ""
        with open(chatPath, 'r', encoding='utf-8') as f:
            for line in f.readlines():
                line = line.strip()
                if ',' in line:
                    datetime = line.split(',')[0].split(' ')
                    if len(datetime) >= 3:
                        y = re.sub(r'[^0-9]', '', datetime[0])
                        m = re.sub(r'[^0-9]', '', datetime[1])
                        d = re.sub(r'[^0-9]', '', datetime[2])
                        cur = y + "-" + m + "-" + d

                if '.jpg' in line \
                        or '.jpeg' in line\
                        or '.JPG' in line\
                        or '.JPEG' in line:
                    if cur not in res.keys():
                        res[cur] = []
                    res[cur].append(line.split(' ')[-1])
        return res

    @Slot(result=bool)
    def check(self):
        return self.isRunning()

    @Slot(str, dict)
    def startMake(self, path, imgNames):
        self.path = path
        self.imgNames = imgNames
        self.start()